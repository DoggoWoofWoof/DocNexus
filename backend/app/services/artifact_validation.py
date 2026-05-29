from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from sqlmodel import Session

from backend.app.db.models import Artifact
from backend.app.schemas.artifact import ArtifactType, ArtifactValidationCheck, ArtifactValidationResult
from backend.app.services.artifacts import get_artifact


REQUIRED_REPORT_SECTIONS = [
    "Executive Summary",
    "Physician Landscape Overview",
    "Geographic & Specialty Distribution",
    "Key Insights & Implications",
    "Recommended Next Steps",
]

REQUIRED_EXCEL_SHEETS = [
    "Raw Physician Data",
    "State x Specialty Summary",
    "ICD-10 Breakdown",
]


def validate_artifacts(
    session: Session,
    artifact_ids: list[str],
) -> list[ArtifactValidationResult]:
    results: list[ArtifactValidationResult] = []
    seen: set[str] = set()
    for artifact_id in artifact_ids:
        if artifact_id in seen:
            continue
        seen.add(artifact_id)
        artifact = get_artifact(session, artifact_id)
        results.append(validate_artifact(artifact))
    return results


def validate_artifact(artifact: Artifact) -> ArtifactValidationResult:
    artifact_type = ArtifactType(artifact.type)
    path = Path(artifact.local_path)

    if artifact_type == ArtifactType.pptx:
        checks = _validate_pptx(path)
    elif artifact_type == ArtifactType.xlsx:
        checks = _validate_xlsx(path)
    elif artifact_type == ArtifactType.markdown:
        checks = _validate_markdown(path)
    elif artifact_type in {ArtifactType.chart_png, ArtifactType.chart_svg}:
        checks = _validate_chart(path, artifact_type)
    else:
        checks = [_check("known_type", True, f"No specialized validator required for {artifact_type.value}.")]

    passed_count = sum(1 for check in checks if check.passed)
    score = round((passed_count / len(checks)) * 100) if checks else 0
    return ArtifactValidationResult(
        artifact_id=artifact.id,
        artifact_type=artifact_type,
        source_agent=artifact.source_agent,
        passed=all(check.passed for check in checks),
        score=score,
        checks=checks,
    )


def _validate_pptx(path: Path) -> list[ArtifactValidationCheck]:
    if not path.exists():
        return [_check("file_exists", False, "PPTX file does not exist.")]

    deck = Presentation(path)
    slide_text = [_slide_text(slide) for slide in deck.slides]
    combined = "\n".join(slide_text).lower()
    return [
        _check("file_exists", True, "PPTX file exists.", {"path": str(path)}),
        _check("slide_count", len(deck.slides) >= 4, "PPTX has at least four slides.", {"slideCount": len(deck.slides)}),
        _check("title_slide", bool(slide_text and slide_text[0].strip()), "Title slide contains text."),
        _check("overview_slide", "overview" in combined, "Deck includes a population overview slide."),
        _check("insight_slide", "insight" in combined, "Deck includes an insights slide."),
        _check("top_physicians", "top 10" in combined or "top physicians" in combined, "Deck includes a top physicians table slide."),
    ]


def _validate_xlsx(path: Path) -> list[ArtifactValidationCheck]:
    if not path.exists():
        return [_check("file_exists", False, "XLSX file does not exist.")]

    workbook = load_workbook(path, read_only=True, data_only=True)
    sheet_names = workbook.sheetnames
    checks = [
        _check("file_exists", True, "XLSX file exists.", {"path": str(path)}),
        _check("required_sheets", all(name in sheet_names for name in REQUIRED_EXCEL_SHEETS), "Workbook contains required sheets.", {"sheetNames": sheet_names}),
    ]
    for sheet_name in REQUIRED_EXCEL_SHEETS:
        if sheet_name in workbook:
            sheet = workbook[sheet_name]
            checks.append(
                _check(
                    f"{sheet_name.lower().replace(' ', '_')}_rows",
                    sheet.max_row > 1,
                    f"{sheet_name} contains data rows.",
                    {"rowCount": sheet.max_row},
                )
            )
    workbook.close()
    return checks


def _validate_markdown(path: Path) -> list[ArtifactValidationCheck]:
    if not path.exists():
        return [_check("file_exists", False, "Markdown report does not exist.")]

    text = path.read_text(encoding="utf-8")
    normalized = text.lower()
    checks = [
        _check("file_exists", True, "Markdown report exists.", {"path": str(path)}),
        _check("non_empty", len(text.strip()) > 100, "Markdown report has substantive content.", {"characterCount": len(text)}),
    ]
    for section in REQUIRED_REPORT_SECTIONS:
        checks.append(
            _check(
                f"section_{section.lower().replace(' ', '_').replace('&', 'and')}",
                section.lower() in normalized,
                f"Report includes {section}.",
            )
        )
    return checks


def _validate_chart(path: Path, artifact_type: ArtifactType) -> list[ArtifactValidationCheck]:
    exists = path.exists()
    checks = [_check("file_exists", exists, "Chart file exists.", {"path": str(path)})]
    if exists:
        size = path.stat().st_size
        checks.append(_check("non_empty", size > 0, "Chart file is non-empty.", {"fileSizeBytes": size}))
        if artifact_type == ArtifactType.chart_png:
            checks.append(_check("png_signature", path.read_bytes().startswith(b"\x89PNG"), "PNG chart has a valid signature."))
    return checks


def _slide_text(slide) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            values.append(shape.text)
    return "\n".join(values)


def _check(
    name: str,
    passed: bool,
    message: str,
    metadata: dict[str, object] | None = None,
) -> ArtifactValidationCheck:
    return ArtifactValidationCheck(name=name, passed=passed, message=message, metadata=metadata or {})
